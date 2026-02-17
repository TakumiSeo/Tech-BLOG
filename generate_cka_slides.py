#!/usr/bin/env python3
"""
CKA学習メモ → PowerPoint 自動生成スクリプト
デザイン: Azure Blue ベース × プレミアムミニマル（ブログテーマ準拠）
"""

import base64
import io
from pathlib import Path

import requests
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

# ──────────────────────────────────────────────
#  Design Tokens (ブログテーマ準拠)
# ──────────────────────────────────────────────
AZURE_BLUE = RGBColor(0x00, 0x78, 0xD4)
AZURE_LIGHT = RGBColor(0x00, 0xBC, 0xF2)
AZURE_DARK = RGBColor(0x00, 0x5A, 0x9E)
HERO_DARK = RGBColor(0x00, 0x14, 0x29)
HERO_MID = RGBColor(0x00, 0x3D, 0x5C)
ACCENT_GOLD = RGBColor(0xFF, 0xB9, 0x00)
BG_PRIMARY = RGBColor(0xF7, 0xF7, 0xFB)
TEXT_PRIMARY = RGBColor(0x0F, 0x17, 0x2A)
TEXT_MUTED = RGBColor(0x4B, 0x55, 0x63)
CODE_BG = RGBColor(0x0F, 0x17, 0x2A)
CODE_TEXT = RGBColor(0xF8, 0xFA, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT_BODY = "Meiryo"
FONT_CODE = "Consolas"

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# ──────────────────────────────────────────────
#  Mermaid Diagrams
# ──────────────────────────────────────────────
MERMAID_FLOW = """flowchart TB
    A[OS準備<br/>swap off / sysctl / kernel modules] --> B[コンテナランタイム<br/>containerd + SystemdCgroup]
    B --> C[Kubernetesパッケージ<br/>kubeadm/kubelet/kubectl]
    C --> D[kubeadm設定<br/>controlPlaneEndpoint / podSubnet]
    D --> E[kubeadm init<br/>クラスタ初期化 + 証明書準備]
    E --> F[kubeconfig設定<br/>kubectlが使える状態にする]
    F --> G[CNI適用<br/>例: Calico/Flannel/Cilium]
    G --> H[ノード追加<br/>kubeadm join]"""

MERMAID_CLUSTER = """flowchart LR
    subgraph CP[Control Plane Node]
        direction TB
        API[kube-apiserver]
        ETCD[etcd]
        CTRL[controller-manager]
        SCH[scheduler]
    end
    subgraph W1[Worker Node]
        direction TB
        K1[kubelet] --> R1[containerd]
        P1[Pods]
    end
    subgraph W2[Worker Node]
        direction TB
        K2[kubelet] --> R2[containerd]
        P2[Pods]
    end
    API --- K1
    API --- K2"""


# ──────────────────────────────────────────────
#  Helper Functions
# ──────────────────────────────────────────────
def fetch_mermaid_png(mermaid_code: str) -> bytes | None:
    """mermaid.ink API 経由で Mermaid 図を PNG として取得する。"""
    try:
        encoded = base64.urlsafe_b64encode(mermaid_code.encode("utf-8")).decode("ascii")
        url = f"https://mermaid.ink/img/{encoded}?theme=neutral&bgColor=F7F7FB"
        resp = requests.get(url, timeout=30)
        resp.raise_for_status()
        return resp.content
    except Exception as e:
        print(f"  ⚠ Mermaid 図の取得に失敗しました: {e}")
        return None


def set_slide_bg_solid(slide, color: RGBColor):
    """スライド背景を単色に設定。"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def set_slide_bg_gradient(slide, color1: RGBColor, color2: RGBColor, color3: RGBColor):
    """スライド背景を3色グラデーションに設定（135°）。"""
    bg = slide.background
    fill = bg.fill
    fill.gradient()
    fill.gradient_angle = 135.0
    stops = fill.gradient_stops
    stops[0].color.rgb = color1
    stops[0].position = 0.0
    stops[1].color.rgb = color2
    stops[1].position = 0.5
    # 3つ目の stop を追加
    stop_el = stops[1]._element
    new_stop = stop_el.makeelement(qn("a:gs"), {"pos": "100000"})
    srgb_val = str(color3)  # RGBColor.__str__ returns "RRGGBB"
    solid = new_stop.makeelement(qn("a:srgbClr"), {"val": srgb_val})
    new_stop.append(solid)
    stop_el.getparent().append(new_stop)


def add_textbox(slide, left, top, width, height, text, font_name=FONT_BODY,
                font_size=Pt(14), font_color=TEXT_PRIMARY, bold=False,
                alignment=PP_ALIGN.LEFT, line_spacing=1.5):
    """テキストボックスを追加するユーティリティ。"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = font_size
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.alignment = alignment
    p.line_spacing = line_spacing
    return txBox, tf


def add_bullet_list(tf, items, font_name=FONT_BODY, font_size=Pt(14),
                    font_color=TEXT_PRIMARY, bullet_color=AZURE_BLUE,
                    line_spacing=1.6, space_after=Pt(8)):
    """テキストフレームに箇条書きを追加。"""
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.name = font_name
        p.font.size = font_size
        p.font.color.rgb = font_color
        p.line_spacing = line_spacing
        p.space_after = space_after
        # ブレットの色を設定
        pPr = p._element.get_or_add_pPr()
        buClr = pPr.makeelement(qn("a:buClr"), {})
        srgb = buClr.makeelement(
            qn("a:srgbClr"),
            {"val": f"{bullet_color.red:02X}{bullet_color.green:02X}{bullet_color.blue:02X}"},
        )
        buClr.append(srgb)
        pPr.append(buClr)
        buChar = pPr.makeelement(qn("a:buChar"), {"char": "●"})
        pPr.append(buChar)


def add_code_block(slide, left, top, width, height, code_text):
    """角丸ダーク背景のコードブロックを追加。"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.adjustments[0] = 0.03
    shape.fill.solid()
    shape.fill.fore_color.rgb = CODE_BG
    shape.line.fill.background()  # 枠線なし

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.3)
    tf.margin_right = Inches(0.3)
    tf.margin_top = Inches(0.2)
    tf.margin_bottom = Inches(0.2)
    tf.vertical_anchor = MSO_ANCHOR.TOP

    for i, line in enumerate(code_text.strip().split("\n")):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.name = FONT_CODE
        p.font.size = Pt(12)
        p.font.color.rgb = CODE_TEXT
        p.line_spacing = 1.5
        p.space_after = Pt(2)
    return shape


def add_slide_number(slide, number):
    """右下にスライド番号を追加。"""
    txBox = slide.shapes.add_textbox(
        SLIDE_WIDTH - Inches(0.8), SLIDE_HEIGHT - Inches(0.5),
        Inches(0.6), Inches(0.3),
    )
    p = txBox.text_frame.paragraphs[0]
    p.text = str(number)
    p.font.name = FONT_BODY
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_MUTED
    p.alignment = PP_ALIGN.RIGHT


def add_accent_bar(slide, left, top, width, height, color=AZURE_BLUE):
    """装飾用アクセントバー（角丸長方形）を追加。"""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    bar.adjustments[0] = 0.5
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return bar


def add_section_title(slide, text, subtitle=None):
    """スライド上部にセクションタイトルを追加（Azure Blue アクセントバー付き）。"""
    add_accent_bar(slide, Inches(0.6), Inches(0.5), Inches(0.08), Inches(0.55))
    add_textbox(
        slide, Inches(0.9), Inches(0.4), Inches(10), Inches(0.7),
        text, font_size=Pt(28), font_color=TEXT_PRIMARY, bold=True,
    )
    if subtitle:
        add_textbox(
            slide, Inches(0.9), Inches(1.0), Inches(10), Inches(0.4),
            subtitle, font_size=Pt(14), font_color=TEXT_MUTED,
        )


def add_numbered_item(slide, number, title, description, left, top, width):
    """番号付きステップアイテムを追加（番号は Azure Blue 丸）。"""
    # 番号の丸
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left, top, Inches(0.45), Inches(0.45)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = AZURE_BLUE
    circle.line.fill.background()
    tf_c = circle.text_frame
    tf_c.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf_c.paragraphs[0]
    p.text = str(number)
    p.font.name = FONT_BODY
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # タイトル
    add_textbox(
        slide, left + Inches(0.6), top - Inches(0.02),
        width - Inches(0.6), Inches(0.35),
        title, font_size=Pt(16), font_color=TEXT_PRIMARY, bold=True,
    )
    # 説明
    add_textbox(
        slide, left + Inches(0.6), top + Inches(0.35),
        width - Inches(0.6), Inches(0.7),
        description, font_size=Pt(12), font_color=TEXT_MUTED, line_spacing=1.4,
    )


def add_warning_card(slide, text, left, top, width, height):
    """⚠ 注意カード（ゴールドアクセント角丸カード）を追加。"""
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    card.adjustments[0] = 0.05
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = ACCENT_GOLD
    card.line.width = Pt(1.5)

    # ゴールド左バー
    add_accent_bar(
        slide, left + Inches(0.05), top + Inches(0.1),
        Inches(0.06), height - Inches(0.2), ACCENT_GOLD,
    )

    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.3)
    tf.margin_top = Inches(0.15)
    tf.margin_bottom = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = f"⚠  {text}"
    p.font.name = FONT_BODY
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT_PRIMARY
    p.line_spacing = 1.4
    return card


def add_shadow_to_shape(shape):
    """シェイプにソフトシャドウを追加。"""
    spPr = shape._element.spPr
    effectLst = spPr.makeelement(qn("a:effectLst"), {})
    outerShdw = effectLst.makeelement(qn("a:outerShdw"), {
        "blurRad": "76200", "dist": "38100", "dir": "5400000",
        "rotWithShape": "0",
    })
    srgbClr = outerShdw.makeelement(qn("a:srgbClr"), {"val": "0F172A"})
    alpha = srgbClr.makeelement(qn("a:alpha"), {"val": "8000"})
    srgbClr.append(alpha)
    outerShdw.append(srgbClr)
    effectLst.append(outerShdw)
    spPr.append(effectLst)


# ──────────────────────────────────────────────
#  Slide Builders
# ──────────────────────────────────────────────
def build_slide_1_title(prs):
    """スライド1: タイトル"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    set_slide_bg_gradient(slide, HERO_DARK, HERO_MID, AZURE_DARK)

    # メインタイトル
    add_textbox(
        slide, Inches(1.2), Inches(1.8), Inches(10.5), Inches(1.2),
        "CKA学習メモ", font_size=Pt(44), font_color=WHITE, bold=True,
        alignment=PP_ALIGN.LEFT,
    )
    # サブタイトル
    add_textbox(
        slide, Inches(1.2), Inches(3.0), Inches(10.5), Inches(0.8),
        "kubeadm 演習メモ（課題3.1）", font_size=Pt(24),
        font_color=AZURE_LIGHT, alignment=PP_ALIGN.LEFT,
    )
    # 説明
    add_textbox(
        slide, Inches(1.2), Inches(4.0), Inches(10.5), Inches(0.6),
        "kubeadm を使ったKubernetesクラスタ構築手順（init/CNI/join）を図解して整理",
        font_size=Pt(14), font_color=RGBColor(0xD2, 0xD7, 0xE1),
        alignment=PP_ALIGN.LEFT,
    )

    # タグバッジ
    tags = ["kubernetes", "cka", "cni"]
    tag_left = Inches(1.2)
    for tag in tags:
        badge = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            tag_left, Inches(4.8), Inches(1.4), Inches(0.38),
        )
        badge.adjustments[0] = 0.5
        badge.fill.solid()
        badge.fill.fore_color.rgb = ACCENT_GOLD
        badge.line.fill.background()
        tf = badge.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = tag
        p.font.name = FONT_BODY
        p.font.size = Pt(11)
        p.font.color.rgb = HERO_DARK
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        tag_left += Inches(1.6)

    # 日付
    add_textbox(
        slide, Inches(1.2), Inches(5.8), Inches(4), Inches(0.4),
        "2026-02-14", font_size=Pt(12),
        font_color=RGBColor(0x8B, 0x95, 0xA5), alignment=PP_ALIGN.LEFT,
    )

    # 装飾ライン
    add_accent_bar(slide, Inches(1.2), Inches(3.7), Inches(3), Inches(0.04), AZURE_BLUE)

    add_slide_number(slide, 1)


def build_slide_2_agenda(prs):
    """スライド2: アジェンダ"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg_solid(slide, BG_PRIMARY)
    add_section_title(slide, "アジェンダ")

    agenda_items = [
        ("01", "全体フロー図", "kubeadm クラスタ構築の全体像をビジュアルで確認"),
        ("02", "ステップ別解説", "OS準備からノード参加まで、8ステップを詳細に解説"),
        ("03", "注意点（詰まりどころ）", "ファイアウォール、CNI不一致など、よくあるハマりポイント"),
    ]

    for i, (num, title, desc) in enumerate(agenda_items):
        card_top = Inches(1.8) + Inches(i * 1.7)
        # カード背景
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(1.0), card_top, Inches(11), Inches(1.4),
        )
        card.adjustments[0] = 0.06
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.fill.background()
        add_shadow_to_shape(card)

        # 番号
        add_textbox(
            slide, Inches(1.4), card_top + Inches(0.3), Inches(0.8), Inches(0.6),
            num, font_size=Pt(32), font_color=AZURE_BLUE, bold=True,
        )
        # アクセントバー
        add_accent_bar(
            slide, Inches(2.3), card_top + Inches(0.25),
            Inches(0.04), Inches(0.9), AZURE_BLUE,
        )
        # タイトル
        add_textbox(
            slide, Inches(2.7), card_top + Inches(0.25), Inches(8), Inches(0.45),
            title, font_size=Pt(20), font_color=TEXT_PRIMARY, bold=True,
        )
        # 説明
        add_textbox(
            slide, Inches(2.7), card_top + Inches(0.7), Inches(8), Inches(0.4),
            desc, font_size=Pt(13), font_color=TEXT_MUTED,
        )

    add_slide_number(slide, 2)


def build_slide_3_flow(prs, flow_png: bytes | None):
    """スライド3: 全体フロー図"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg_solid(slide, BG_PRIMARY)
    add_section_title(slide, "全体フロー図", "kubeadm クラスタ構築の流れ")

    if flow_png:
        img_stream = io.BytesIO(flow_png)
        # フレーム（白背景カード）
        frame = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(1.2), Inches(1.5), Inches(10.8), Inches(5.6),
        )
        frame.adjustments[0] = 0.02
        frame.fill.solid()
        frame.fill.fore_color.rgb = WHITE
        frame.line.color.rgb = RGBColor(0xE5, 0xE7, 0xEB)
        frame.line.width = Pt(1)
        add_shadow_to_shape(frame)

        slide.shapes.add_picture(
            img_stream, Inches(1.5), Inches(1.7), Inches(10.2), Inches(5.2),
        )
    else:
        # フォールバック: テキストで表示
        steps = (
            "OS準備 → コンテナランタイム → Kubernetesパッケージ → "
            "kubeadm設定 → kubeadm init → kubeconfig設定 → CNI適用 → ノード追加"
        )
        add_textbox(
            slide, Inches(1.5), Inches(3.0), Inches(10), Inches(1.5),
            steps, font_size=Pt(18), font_color=TEXT_PRIMARY,
            alignment=PP_ALIGN.CENTER,
        )

    add_slide_number(slide, 3)


def build_slide_4_cluster(prs, cluster_png: bytes | None):
    """スライド4: クラスタ構成イメージ"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg_solid(slide, BG_PRIMARY)
    add_section_title(slide, "クラスタ構成イメージ", "Control Plane + Worker Nodes")

    if cluster_png:
        img_stream = io.BytesIO(cluster_png)
        frame = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(1.5), Inches(1.5), Inches(10.3), Inches(5.6),
        )
        frame.adjustments[0] = 0.02
        frame.fill.solid()
        frame.fill.fore_color.rgb = WHITE
        frame.line.color.rgb = RGBColor(0xE5, 0xE7, 0xEB)
        frame.line.width = Pt(1)
        add_shadow_to_shape(frame)

        slide.shapes.add_picture(
            img_stream, Inches(1.8), Inches(1.7), Inches(9.7), Inches(5.2),
        )
    else:
        desc = (
            "Control Plane Node: kube-apiserver / etcd / controller-manager / scheduler\n"
            "Worker Node 1: kubelet → containerd → Pods\n"
            "Worker Node 2: kubelet → containerd → Pods"
        )
        add_textbox(
            slide, Inches(2.0), Inches(3.0), Inches(9), Inches(2.0),
            desc, font_size=Pt(16), font_color=TEXT_PRIMARY,
            alignment=PP_ALIGN.CENTER,
        )

    add_slide_number(slide, 4)


def build_slide_5_steps_1(prs):
    """スライド5: ステップ解説① (1-4)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg_solid(slide, BG_PRIMARY)
    add_section_title(slide, "ステップ別解説 ①", "OS準備 〜 kubeadm-config.yaml")

    steps = [
        (
            "1", "OS/カーネル準備",
            "swapoff -a でスワップ無効化。modprobe overlay / br_netfilter で"
            "カーネルモジュール読込。sysctl で ip_forward=1 を設定。",
        ),
        (
            "2", "コンテナランタイム導入",
            "containerd をインストール。SystemdCgroup = true に設定して"
            " kubelet と cgroup ドライバを統一。",
        ),
        (
            "3", "Kubernetes パッケージ導入",
            "kubeadm（init/join担当）、kubelet（Pod実行エージェント）、"
            "kubectl（操作CLI）をインストール。",
        ),
        (
            "4", "kubeadm-config.yaml 作成",
            "controlPlaneEndpoint に名前を使用（証明書の耐久性向上）。"
            "podSubnet を CNI 設定と一致させる。",
        ),
    ]

    col_left = [Inches(0.8), Inches(6.8)]
    for i, (num, title, desc) in enumerate(steps):
        col = i % 2
        row = i // 2
        left = col_left[col]
        top = Inches(1.7) + Inches(row * 2.6)
        add_numbered_item(slide, num, title, desc, left, top, Inches(5.5))

    add_slide_number(slide, 5)


def build_slide_6_steps_2(prs):
    """スライド6: ステップ解説② (5-8)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg_solid(slide, BG_PRIMARY)
    add_section_title(slide, "ステップ別解説 ②", "kubeadm init 〜 ノード参加")

    steps = [
        (
            "5", "kubeadm init 実行",
            "Control Plane コンポーネントを起動。join コマンドが出力される。"
            "--upload-certs で CP追加を容易にする。",
        ),
        (
            "6", "kubeconfig 設定",
            "admin.conf を ~/.kube/config にコピー。"
            "これで kubectl が管理者権限で操作可能になる。",
        ),
        (
            "7", "CNI（Podネットワーク）適用",
            "CNI が入って初めて Pod に IP が付与される。"
            "ここが終わると Node が Ready に遷移する。",
        ),
        (
            "8", "ノード追加（kubeadm join）",
            "token / certificate-key は機密情報。Git には実値を残さない。"
            "証明書は2時間で自動削除される。",
        ),
    ]

    col_left = [Inches(0.8), Inches(6.8)]
    for i, (num, title, desc) in enumerate(steps):
        col = i % 2
        row = i // 2
        left = col_left[col]
        top = Inches(1.7) + Inches(row * 2.6)
        add_numbered_item(slide, num, title, desc, left, top, Inches(5.5))

    add_slide_number(slide, 6)


def build_slide_7_pitfalls(prs):
    """スライド7: 注意点"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg_solid(slide, BG_PRIMARY)
    add_section_title(slide, "注意点（詰まりどころ）", "クラスタ構築でよくある落とし穴")

    warnings = [
        "ファイアウォール / セキュリティグループ: ノード間通信や API（6443）が塞がれていると詰まりやすい",
        "CNI と podSubnet の不一致: Node が NotReady のまま / Pod が通信できない原因になりやすい",
        "YAML はタブ禁止: インデントはスペースのみ。コピペ時にタブが混入しがち",
        "join の情報は秘匿: --token / --certificate-key は共有しない。Git 管理のメモには実値を残さない",
    ]

    for i, text in enumerate(warnings):
        col = i % 2
        row = i // 2
        left = Inches(0.8) + Inches(col * 6.2)
        top = Inches(1.8) + Inches(row * 2.5)
        add_warning_card(slide, text, left, top, Inches(5.8), Inches(2.0))

    add_slide_number(slide, 7)


def build_slide_8_commands(prs):
    """スライド8: 主要コマンド一覧"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg_solid(slide, BG_PRIMARY)
    add_section_title(slide, "主要コマンド一覧", "クラスタ構築で使用する主なコマンド")

    commands = [
        (
            "kubeadm init（ログ保存）",
            "kubeadm init --config=kubeadm-config.yaml \\\n"
            "  --upload-certs | tee kubeadm-init.out",
        ),
        (
            "kubeconfig 設定",
            "mkdir -p $HOME/.kube\n"
            "sudo cp -i /etc/kubernetes/admin.conf \\\n"
            "  $HOME/.kube/config\n"
            "sudo chown $(id -u):$(id -g) $HOME/.kube/config",
        ),
        (
            "CNI デプロイ",
            "kubectl apply -f <podnetwork>.yaml",
        ),
        (
            "Worker ノード参加",
            "kubeadm join k8scp:6443 --token <token> \\\n"
            "  --discovery-token-ca-cert-hash sha256:<hash>",
        ),
    ]

    for i, (title, code) in enumerate(commands):
        col = i % 2
        row = i // 2
        left = Inches(0.6) + Inches(col * 6.2)
        top = Inches(1.7) + Inches(row * 2.8)

        # ラベル
        add_textbox(
            slide, left, top, Inches(5.8), Inches(0.35),
            title, font_size=Pt(13), font_color=AZURE_BLUE, bold=True,
        )
        # コードブロック
        add_code_block(slide, left, top + Inches(0.4), Inches(5.8), Inches(2.0), code)

    add_slide_number(slide, 8)


def build_slide_9_summary(prs):
    """スライド9: まとめ"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg_gradient(slide, HERO_DARK, HERO_MID, AZURE_DARK)

    add_textbox(
        slide, Inches(1.2), Inches(1.5), Inches(10.5), Inches(0.8),
        "まとめ", font_size=Pt(36), font_color=WHITE, bold=True,
    )
    add_accent_bar(slide, Inches(1.2), Inches(2.3), Inches(3), Inches(0.04), AZURE_BLUE)

    points = [
        "kubeadm を使ったクラスタ構築は「OS準備 → ランタイム → K8s → init → CNI → join」の流れ",
        "CNI の podSubnet 設定を kubeadm-config.yaml と一致させることが重要",
        "join に使う token / certificate-key は機密情報として取り扱う",
        "ファイアウォール設定と YAML のインデント（スペースのみ）に注意",
    ]

    for i, point in enumerate(points):
        top = Inches(2.8) + Inches(i * 0.85)
        # チェックマーク
        add_textbox(
            slide, Inches(1.2), top, Inches(0.4), Inches(0.35),
            "✓", font_size=Pt(18), font_color=ACCENT_GOLD, bold=True,
        )
        add_textbox(
            slide, Inches(1.7), top, Inches(10), Inches(0.6),
            point, font_size=Pt(16),
            font_color=RGBColor(0xD2, 0xD7, 0xE1),
        )

    # 参考リンク
    add_textbox(
        slide, Inches(1.2), Inches(6.2), Inches(10), Inches(0.4),
        "参考: https://kubernetes.io/docs/concepts/cluster-administration/addons/",
        font_size=Pt(11), font_color=RGBColor(0x8B, 0x95, 0xA5),
    )

    add_slide_number(slide, 9)


# ──────────────────────────────────────────────
#  Main
# ──────────────────────────────────────────────
def main():
    print("📊 CKA学習メモ PowerPoint 生成中...")

    # Mermaid 図を取得
    print("  → Mermaid フロー図を取得中...")
    flow_png = fetch_mermaid_png(MERMAID_FLOW)
    print("  → Mermaid クラスタ構成図を取得中...")
    cluster_png = fetch_mermaid_png(MERMAID_CLUSTER)

    # Presentation 作成（16:9）
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # 全9スライドを構築
    print("  → スライドを構築中...")
    build_slide_1_title(prs)
    build_slide_2_agenda(prs)
    build_slide_3_flow(prs, flow_png)
    build_slide_4_cluster(prs, cluster_png)
    build_slide_5_steps_1(prs)
    build_slide_6_steps_2(prs)
    build_slide_7_pitfalls(prs)
    build_slide_8_commands(prs)
    build_slide_9_summary(prs)

    # 保存
    output_path = Path(__file__).parent / "cka-training-memo.pptx"
    prs.save(str(output_path))
    print(f"✅ 完了: {output_path}")


if __name__ == "__main__":
    main()
