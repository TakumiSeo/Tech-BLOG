#!/usr/bin/env python
"""Build a 16:9 PowerPoint deck from a JSON slide spec.

Repo-specific helper for the `presentation-deck` skill. Uses python-pptx and
reuses existing draw.io PNG diagrams under content/images/<slug>/.

Usage:
    python build_deck.py <spec.json> [-o output.pptx]

The spec schema is documented in SKILL.md. This script is intentionally
self-contained (stdlib json + python-pptx + Pillow) so it runs in the repo venv
without extra installs.
"""
from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

# 16:9 canvas
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Built-in palettes: (dark, light, accent, muted-on-dark, muted-on-light)
PALETTES = {
    "midnight": {
        "dark": "1E2761", "light": "FFFFFF", "accent": "5B8DEF",
        "ink": "1A1A2E", "sub": "5A6072", "band": "EAF0FB",
    },
    "teal": {
        "dark": "023E45", "light": "FFFFFF", "accent": "02C39A",
        "ink": "10322F", "sub": "4F6B68", "band": "E3F5F0",
    },
    "charcoal": {
        "dark": "262B33", "light": "FFFFFF", "accent": "F2A900",
        "ink": "23272E", "sub": "5A6068", "band": "F1F2F4",
    },
    "berry": {
        "dark": "4A1C40", "light": "FFFFFF", "accent": "E0529C",
        "ink": "2E1329", "sub": "6E5366", "band": "F7E9F2",
    },
}

DEFAULT_FONT = "Yu Gothic UI"  # good Japanese + Latin coverage on Windows


def hex_color(value: str) -> RGBColor:
    return RGBColor.from_string(value.lstrip("#"))


def set_run(run, *, font=DEFAULT_FONT, size=None, bold=None, italic=None, color=None):
    """Set font on a run including the East Asian (ea) typeface for Japanese."""
    if size is not None:
        run.font.size = Pt(size)
    if bold is not None:
        run.font.bold = bold
    if italic is not None:
        run.font.italic = italic
    if color is not None:
        run.font.color.rgb = hex_color(color)
    run.font.name = font  # sets latin
    rpr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rpr.find(qn(tag))
        if el is None:
            el = rpr.makeelement(qn(tag), {})
            rpr.append(el)
        el.set("typeface", font)


def add_rect(slide, x, y, w, h, fill_hex, line_hex=None):
    from pptx.enum.shapes import MSO_SHAPE

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_color(fill_hex)
    if line_hex is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = hex_color(line_hex)
    shape.shadow.inherit = False
    return shape


def add_text(slide, x, y, w, h, *, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.paragraphs[0].alignment = align
    return tf


def add_paragraph(tf, text, *, first=False, **run_kw):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    run = p.add_run()
    run.text = text
    set_run(run, **run_kw)
    return p


def footer(slide, pal, note):
    if not note:
        return
    tf = add_text(slide, Inches(0.5), Inches(7.02), Inches(12.3), Inches(0.4),
                  anchor=MSO_ANCHOR.MIDDLE)
    add_paragraph(tf, note, first=True, font=DEFAULT_FONT, size=10, color=pal["sub"])


def notes(slide, text):
    if text:
        slide.notes_slide.notes_text_frame.text = text


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _resolve_image(image: str, base: str) -> Path:
    """Resolve an image path. Absolute wins; else try CWD (repo root) then spec dir."""
    p = Path(image)
    if p.is_absolute():
        return p
    cwd_rel = (Path.cwd() / p).resolve()
    if cwd_rel.exists():
        return cwd_rel
    spec_rel = (Path(base) / p).resolve()
    if spec_rel.exists():
        return spec_rel
    raise FileNotFoundError(
        f"Image not found: {image!r} (tried {cwd_rel} and {spec_rel})"
    )



def slide_title(prs, pal, s):
    slide = blank(prs)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, pal["dark"])
    add_rect(slide, Inches(0.0), Inches(3.05), Inches(2.2), Inches(0.12), pal["accent"])
    tf = add_text(slide, Inches(0.9), Inches(2.3), Inches(11.5), Inches(2.6),
                  anchor=MSO_ANCHOR.TOP)
    add_paragraph(tf, s["title"], first=True, size=44, bold=True, color=pal["light"])
    if s.get("subtitle"):
        add_paragraph(tf, s["subtitle"], size=24, color=pal["accent"])
    meta = " ｜ ".join(x for x in [s.get("author"), s.get("date")] if x)
    if meta:
        mtf = add_text(slide, Inches(0.9), Inches(6.4), Inches(11.5), Inches(0.6))
        add_paragraph(mtf, meta, first=True, size=14, color=pal["light"])
    notes(slide, s.get("notes"))
    return slide


def slide_section(prs, pal, s):
    slide = blank(prs)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, pal["dark"])
    add_rect(slide, Inches(0.9), Inches(3.5), Inches(1.6), Inches(0.1), pal["accent"])
    tf = add_text(slide, Inches(0.9), Inches(2.7), Inches(11.5), Inches(2.0),
                  anchor=MSO_ANCHOR.TOP)
    if s.get("eyebrow"):
        add_paragraph(tf, s["eyebrow"], first=True, size=16, bold=True, color=pal["accent"])
        add_paragraph(tf, s["title"], size=40, bold=True, color=pal["light"])
    else:
        add_paragraph(tf, s["title"], first=True, size=40, bold=True, color=pal["light"])
    notes(slide, s.get("notes"))
    return slide


def _content_header(slide, pal, title):
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.25), pal["light"])
    add_rect(slide, Inches(0.5), Inches(0.42), Inches(0.16), Inches(0.5), pal["accent"])
    tf = add_text(slide, Inches(0.8), Inches(0.32), Inches(12.0), Inches(0.8),
                  anchor=MSO_ANCHOR.MIDDLE)
    add_paragraph(tf, title, first=True, size=28, bold=True, color=pal["ink"])


def slide_bullets(prs, pal, s):
    slide = blank(prs)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, pal["light"])
    _content_header(slide, pal, s["title"])
    tf = add_text(slide, Inches(0.9), Inches(1.6), Inches(11.6), Inches(5.2))
    for i, item in enumerate(s.get("bullets", [])):
        text, level = (item.get("text"), item.get("level", 0)) if isinstance(item, dict) else (item, 0)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = level
        p.space_after = Pt(10)
        marker = "—  " if level else "●  "
        run = p.add_run(); run.text = marker
        set_run(run, size=20 if not level else 16, bold=True, color=pal["accent"])
        run2 = p.add_run(); run2.text = text
        set_run(run2, size=20 if not level else 16, color=pal["ink"])
    footer(slide, pal, s.get("source"))
    notes(slide, s.get("notes"))
    return slide


def slide_image(prs, pal, s):
    slide = blank(prs)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, pal["light"])
    _content_header(slide, pal, s["title"])
    img_path = _resolve_image(s["image"], s.get("_base", "."))
    avail_x, avail_y = Inches(0.7), Inches(1.5)
    avail_w, avail_h = Inches(11.93), Inches(5.0)
    with Image.open(img_path) as im:
        iw, ih = im.size
    scale = min(avail_w / iw, avail_h / ih)
    w, h = int(iw * scale), int(ih * scale)
    x = avail_x + int((avail_w - w) / 2)
    y = avail_y + int((avail_h - h) / 2)
    slide.shapes.add_picture(str(img_path), Emu(x), Emu(y), width=Emu(w), height=Emu(h))
    if s.get("caption"):
        ctf = add_text(slide, Inches(0.7), Inches(6.55), Inches(11.93), Inches(0.5),
                       align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_paragraph(ctf, s["caption"], first=True, size=12, italic=True, color=pal["sub"])
    footer(slide, pal, s.get("source"))
    notes(slide, s.get("notes"))
    return slide


def slide_quote(prs, pal, s):
    slide = blank(prs)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, pal["dark"])
    qtf = add_text(slide, Inches(1.2), Inches(2.1), Inches(10.9), Inches(3.0),
                   anchor=MSO_ANCHOR.MIDDLE)
    add_paragraph(qtf, "“" + s["text"] + "”", first=True, size=30, bold=True, color=pal["light"])
    if s.get("attribution"):
        atf = add_text(slide, Inches(1.2), Inches(5.0), Inches(10.9), Inches(0.7))
        add_paragraph(atf, "— " + s["attribution"], first=True, size=18, color=pal["accent"])
    footer(slide, pal, s.get("source"))
    notes(slide, s.get("notes"))
    return slide


def slide_closing(prs, pal, s):
    slide = blank(prs)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, pal["dark"])
    add_rect(slide, Inches(0.9), Inches(1.05), Inches(1.6), Inches(0.1), pal["accent"])
    tf = add_text(slide, Inches(0.9), Inches(1.3), Inches(11.5), Inches(1.2))
    add_paragraph(tf, s.get("title", "まとめ"), first=True, size=34, bold=True, color=pal["light"])
    btf = add_text(slide, Inches(0.95), Inches(2.7), Inches(11.4), Inches(3.8))
    for i, item in enumerate(s.get("bullets", [])):
        text = item["text"] if isinstance(item, dict) else item
        p = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
        p.space_after = Pt(10)
        run = p.add_run(); run.text = "✓  "
        set_run(run, size=20, bold=True, color=pal["accent"])
        run2 = p.add_run(); run2.text = text
        set_run(run2, size=20, color=pal["light"])
    footer(slide, pal, s.get("source"))
    notes(slide, s.get("notes"))
    return slide


BUILDERS = {
    "title": slide_title,
    "section": slide_section,
    "bullets": slide_bullets,
    "image": slide_image,
    "quote": slide_quote,
    "closing": slide_closing,
}


def build(spec_path: Path, out_path: Path) -> None:
    spec = json.loads(spec_path.read_text(encoding="utf-8"))
    meta = spec.get("meta", {})
    pal = PALETTES.get(meta.get("theme", "midnight"), PALETTES["midnight"])
    base = spec_path.parent

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    for s in spec.get("slides", []):
        s.setdefault("_base", str(base))
        builder = BUILDERS.get(s.get("type"))
        if builder is None:
            print(f"WARN: unknown slide type {s.get('type')!r}, skipping", file=sys.stderr)
            continue
        builder(prs, pal, s)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    print(f"Wrote {out_path} ({len(prs.slides._sldIdLst)} slides)")


def main(argv=None):
    ap = argparse.ArgumentParser(description="Build a 16:9 PPTX from a JSON slide spec.")
    ap.add_argument("spec", type=Path, help="Path to the JSON slide spec")
    ap.add_argument("-o", "--output", type=Path, help="Output .pptx path")
    args = ap.parse_args(argv)
    out = args.output or Path("decks") / (args.spec.stem.replace(".deck", "") + ".pptx")
    build(args.spec, out)


if __name__ == "__main__":
    main()
